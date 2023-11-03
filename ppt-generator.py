import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Initialize presentation
prs = Presentation()

# Content from the XML
# sections = 

# Get the JSON file path from the user
file_path = input("Please provide the path to the JSON file: ")

# Load the sections from the JSON file
with open(file_path, 'r') as f:
    sections = json.load(f)

# Generate slides from content
for section in sections:
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = section['title']

    if 'description' in section:
        content.text = section['description']
    elif 'points' in section:
        for point in section['points']:
            p = content.text_frame.add_paragraph()
            p.text = point

    if 'code' in section:
        left = Inches(0.5)
        top = Inches(2.5)
        width = Inches(9)
        height = Inches(3)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = section['code']
        p.space_after = Inches(0.1)
        p.line_spacing = 1.25

        # Style the textbox to look like code snippet
        for paragraph in tf.paragraphs:
            for run in paragraph.runs:
                run.font.name = 'Courier New'
                run.font.size = Pt(16)
                run.font.color.rgb = RGBColor(255, 255, 255) # white font
        txBox.fill.solid()
        
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = RGBColor(0, 0, 0) # black background

# Save presentation
prs.save('presentation.pptx')